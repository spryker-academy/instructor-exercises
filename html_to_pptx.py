#!/usr/bin/env python3
"""Convert HTML guides to PowerPoint presentations.

Parses pandoc-generated HTML guides and creates one .pptx per course (basics / intermediate).
Code blocks preserve syntax highlighting from pandoc's HTML spans.
"""

import os
import re
import sys
from pathlib import Path

from bs4 import BeautifulSoup, NavigableString, Tag
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# Spryker brand color
SPRYKER_TEAL = RGBColor(0x00, 0xA8, 0x9B)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
CODE_BG = RGBColor(0xF0, 0xF0, 0xF0)
CODE_COLOR = RGBColor(0x2D, 0x2D, 0x2D)
DIM_COLOR = RGBColor(0x88, 0x88, 0x88)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

CODE_FONT_SIZE = Pt(9)
CODE_FONT_NAME = "Consolas"

# Pandoc syntax highlight class -> color mapping (matches pandoc's CSS)
SYNTAX_COLORS = {
    "kw": RGBColor(0x00, 0x70, 0x20),   # Keyword — green bold
    "cf": RGBColor(0x00, 0x70, 0x20),   # ControlFlow — green bold
    "dt": RGBColor(0x90, 0x20, 0x00),   # DataType — dark red
    "dv": RGBColor(0x40, 0xA0, 0x70),   # DecVal — teal
    "bn": RGBColor(0x40, 0xA0, 0x70),   # BaseN
    "fl": RGBColor(0x40, 0xA0, 0x70),   # Float
    "ch": RGBColor(0x40, 0x70, 0xA0),   # Char — blue
    "st": RGBColor(0x40, 0x70, 0xA0),   # String — blue
    "vs": RGBColor(0x40, 0x70, 0xA0),   # VerbatimString
    "ss": RGBColor(0xBB, 0x66, 0x88),   # SpecialString
    "co": RGBColor(0x60, 0xA0, 0xB0),   # Comment — grey-blue italic
    "an": RGBColor(0x60, 0xA0, 0xB0),   # Annotation
    "cv": RGBColor(0x60, 0xA0, 0xB0),   # CommentVar
    "do": RGBColor(0xBA, 0x21, 0x21),   # Documentation
    "fu": RGBColor(0x06, 0x28, 0x7E),   # Function — dark blue
    "va": RGBColor(0x19, 0x17, 0x7C),   # Variable — purple
    "cn": RGBColor(0x88, 0x00, 0x00),   # Constant — maroon
    "op": RGBColor(0x66, 0x66, 0x66),   # Operator — grey
    "ot": RGBColor(0x00, 0x70, 0x20),   # Other
    "al": RGBColor(0xFF, 0x00, 0x00),   # Alert — red bold
    "er": RGBColor(0xFF, 0x00, 0x00),   # Error — red bold
    "bu": None,                          # BuiltIn — default
    "ex": None,                          # Extension — default
    "im": None,                          # Import — default
    "in": RGBColor(0x60, 0xA0, 0xB0),   # Information
    "pp": RGBColor(0xBC, 0x7A, 0x00),   # Preprocessor
    "sc": RGBColor(0x40, 0x70, 0xA0),   # SpecialChar
    "wa": RGBColor(0x60, 0xA0, 0xB0),   # Warning
}

BOLD_CLASSES = {"kw", "cf", "al", "er", "an", "cv", "in", "wa"}
ITALIC_CLASSES = {"co", "an", "cv", "do", "in", "wa"}

# --- Course objectives ---

BASICS_OBJECTIVES = [
    "Understand Spryker's layered architecture (Yves, Zed, Client, Shared)",
    "Build Back Office pages with Controllers, Twig templates, and Navigation",
    "Define and generate Data Transfer Objects (DTOs) using Transfer XML",
    "Create database schemas with Propel ORM and Schema XML",
    "Implement module layers: Persistence, Business, and Communication",
    "Connect Storefront to Back Office via Client/Stub and BackendGateway",
    "Manage application configuration with Config classes and DI",
    "Extend core Spryker modules following SOLID principles",
]

INTERMEDIATE_OBJECTIVES = [
    "Import structured data using Spryker's Data Import framework",
    "Build advanced Back Office interfaces with tables, forms, and CRUD",
    "Implement Publish & Synchronize for real-time data projection",
    "Integrate Elasticsearch for product and content search",
    "Design RESTful Glue Storefront APIs with resource routing",
    "Model order workflows using the Order Management System (OMS)",
    "Read from Storage (Redis) using the Client layer",
    "Build Yves Storefront pages with Twig, routing, and controllers",
]


def create_presentation():
    """Create a blank widescreen presentation."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def add_title_slide(prs, title, subtitle=""):
    """Add a title slide with Spryker styling."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    left = Inches(1)
    top = Inches(2.2)
    width = Inches(11.333)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = SPRYKER_TEAL
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        top2 = Inches(3.8)
        height2 = Inches(1)
        txBox2 = slide.shapes.add_textbox(left, top2, width, height2)
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = DIM_COLOR
        p2.alignment = PP_ALIGN.CENTER

    # Bottom accent line
    shape = slide.shapes.add_shape(
        1, Inches(4), Inches(5.2), Inches(5.333), Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = SPRYKER_TEAL
    shape.line.fill.background()

    return slide


def add_objectives_slide(prs, objectives):
    """Add a 'What You Will Learn' slide with objectives list."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Heading
    left = Inches(0.7)
    top = Inches(0.4)
    width = Inches(11.9)
    height = Inches(0.9)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "What You Will Learn"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = SPRYKER_TEAL

    # Accent underline
    line_shape = slide.shapes.add_shape(
        1, left, Inches(1.25), Inches(3), Pt(2)
    )
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = SPRYKER_TEAL
    line_shape.line.fill.background()

    # Objectives in two columns
    mid = len(objectives) // 2 + len(objectives) % 2
    col_left = [objectives[:mid], objectives[mid:]]
    col_positions = [Inches(0.9), Inches(6.8)]

    for col_idx, items in enumerate(col_left):
        content_left = col_positions[col_idx]
        content_top = Inches(1.7)
        content_width = Inches(5.5)
        content_height = Inches(5.0)
        txBox2 = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
        tf2 = txBox2.text_frame
        tf2.word_wrap = True

        for i, obj in enumerate(items):
            if i == 0:
                p = tf2.paragraphs[0]
            else:
                p = tf2.add_paragraph()
            run = p.add_run()
            run.text = "\u2713  "
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = SPRYKER_TEAL
            run2 = p.add_run()
            run2.text = obj
            run2.font.size = Pt(15)
            run2.font.color.rgb = DARK_TEXT
            p.space_after = Pt(10)

    return slide


def extract_code_tokens(tag):
    """Extract code tokens with syntax class info from a pandoc sourceCode block.

    Returns a list of lines, where each line is a list of (text, css_class_or_None) tuples.
    """
    code_tag = tag.find("code") or tag
    lines = []
    current_line = []

    def walk(node):
        nonlocal current_line
        if isinstance(node, NavigableString):
            text = str(node)
            # Split by newlines to track line boundaries
            parts = text.split("\n")
            for i, part in enumerate(parts):
                if i > 0:
                    lines.append(current_line)
                    current_line = []
                if part:
                    # Inherit class from parent span
                    css_class = None
                    parent = node.parent
                    if parent and parent.name == "span":
                        classes = parent.get("class", [])
                        for c in classes:
                            if c in SYNTAX_COLORS:
                                css_class = c
                                break
                    current_line.append((part, css_class))
        elif isinstance(node, Tag):
            if node.name == "a":
                return  # skip pandoc line-number anchors
            for child in node.children:
                walk(child)

    walk(code_tag)
    if current_line:
        lines.append(current_line)

    # Strip leading/trailing empty lines
    while lines and not any(text.strip() for text, _ in lines[0]):
        lines.pop(0)
    while lines and not any(text.strip() for text, _ in lines[-1]):
        lines.pop()

    return lines


def add_code_runs_to_paragraph(p, tokens):
    """Add syntax-highlighted runs to a paragraph from a list of (text, css_class) tokens."""
    for text, css_class in tokens:
        run = p.add_run()
        run.text = text
        run.font.size = CODE_FONT_SIZE
        run.font.name = CODE_FONT_NAME

        if css_class and css_class in SYNTAX_COLORS:
            color = SYNTAX_COLORS[css_class]
            if color:
                run.font.color.rgb = color
            else:
                run.font.color.rgb = CODE_COLOR
            if css_class in BOLD_CLASSES:
                run.font.bold = True
            if css_class in ITALIC_CLASSES:
                run.font.italic = True
        else:
            run.font.color.rgb = CODE_COLOR


def add_content_slide(prs, heading, content_blocks):
    """Add a content slide with heading and body content."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Heading
    left = Inches(0.7)
    top = Inches(0.4)
    width = Inches(11.9)
    height = Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = heading
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = SPRYKER_TEAL

    # Accent underline
    line_shape = slide.shapes.add_shape(
        1, left, Inches(1.15), Inches(2), Pt(2)
    )
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = SPRYKER_TEAL
    line_shape.line.fill.background()

    # Content area
    content_top = Inches(1.5)
    content_left = Inches(0.7)
    content_width = Inches(11.9)
    content_height = Inches(5.5)
    txBox2 = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    first = True
    for block in content_blocks:
        btype = block["type"]

        if btype == "code_highlighted":
            # Syntax-highlighted code block
            code_lines = block["lines"]
            for i, tokens in enumerate(code_lines):
                if first:
                    p = tf2.paragraphs[0]
                    first = False
                else:
                    p = tf2.add_paragraph()
                p.space_after = Pt(0)
                p.space_before = Pt(0)
                if tokens:
                    add_code_runs_to_paragraph(p, tokens)
                else:
                    # Empty line — add a space to preserve line height
                    run = p.add_run()
                    run.text = " "
                    run.font.size = CODE_FONT_SIZE
                    run.font.name = CODE_FONT_NAME
            continue

        text = block["text"].strip()
        if not text:
            continue

        if first:
            p = tf2.paragraphs[0]
            first = False
        else:
            p = tf2.add_paragraph()

        if btype == "paragraph":
            p.text = text
            p.font.size = Pt(16)
            p.font.color.rgb = DARK_TEXT
            p.space_after = Pt(8)

        elif btype == "list_item":
            level = block.get("level", 0)
            indent = "    " * level + "\u2022 "
            p.text = indent + text
            p.font.size = Pt(15)
            p.font.color.rgb = DARK_TEXT
            p.level = level
            p.space_after = Pt(4)

        elif btype == "code":
            # Fallback plain code (no syntax spans found)
            lines = text.split("\n")
            for i, line in enumerate(lines):
                if i > 0:
                    p = tf2.add_paragraph()
                p.text = line
                p.font.size = CODE_FONT_SIZE
                p.font.name = CODE_FONT_NAME
                p.font.color.rgb = CODE_COLOR
                p.space_after = Pt(0)

        elif btype == "heading3":
            p.text = text
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = DARK_TEXT
            p.space_before = Pt(12)
            p.space_after = Pt(6)

        elif btype == "heading4":
            p.text = text
            p.font.size = Pt(17)
            p.font.bold = True
            p.font.color.rgb = DARK_TEXT
            p.space_before = Pt(8)
            p.space_after = Pt(4)

        elif btype == "table_text":
            lines = text.split("\n")
            for i, line in enumerate(lines):
                if i > 0:
                    p = tf2.add_paragraph()
                p.text = line
                p.font.size = Pt(10)
                p.font.name = CODE_FONT_NAME
                p.font.color.rgb = DARK_TEXT
                p.space_after = Pt(1)

        elif btype == "blockquote":
            p.text = "\u275D " + text
            p.font.size = Pt(14)
            p.font.italic = True
            p.font.color.rgb = DIM_COLOR
            p.space_after = Pt(8)

    return slide


def extract_text(tag):
    """Extract clean text from a tag, collapsing whitespace."""
    text = tag.get_text(separator=" ", strip=True)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def extract_code_text(tag):
    """Extract code preserving line breaks."""
    text = tag.get_text()
    lines = text.split("\n")
    while lines and not lines[0].strip():
        lines.pop(0)
    while lines and not lines[-1].strip():
        lines.pop()
    return "\n".join(lines)


def table_to_text(table_tag):
    """Convert an HTML table to a simple text representation."""
    rows = []
    for tr in table_tag.find_all("tr"):
        cells = []
        for td in tr.find_all(["th", "td"]):
            cells.append(extract_text(td))
        rows.append(cells)

    if not rows:
        return ""

    col_widths = []
    for row in rows:
        for i, cell in enumerate(row):
            while len(col_widths) <= i:
                col_widths.append(0)
            col_widths[i] = max(col_widths[i], len(cell))

    lines = []
    for ri, row in enumerate(rows):
        parts = []
        for i, cell in enumerate(row):
            w = col_widths[i] if i < len(col_widths) else len(cell)
            parts.append(cell.ljust(w))
        lines.append("  ".join(parts))
        if ri == 0:
            lines.append("  ".join("-" * w for w in col_widths))

    return "\n".join(lines)


def _visual_lines(block):
    """Estimate how many visual lines a block occupies on a slide.

    A widescreen slide content area fits roughly 18-20 code lines at Pt(9),
    or ~14 body-text lines at Pt(15-16), or ~12 heading+body combos.
    We normalize everything to 'code-line equivalents' where 1.0 = one Pt(9) line.
    """
    btype = block["type"]
    if btype == "code_highlighted":
        return len(block["lines"])
    elif btype == "code":
        return block["text"].count("\n") + 1
    elif btype == "table_text":
        return block["text"].count("\n") + 1
    elif btype in ("paragraph", "blockquote"):
        # Paragraphs at Pt(16) ~ 1.8x code line height; long text wraps
        text = block.get("text", "")
        wrapped = max(1, len(text) // 100 + 1)  # rough wrap estimate
        return wrapped * 1.8
    elif btype == "list_item":
        text = block.get("text", "")
        wrapped = max(1, len(text) // 100 + 1)
        return wrapped * 1.6
    elif btype == "heading3":
        return 2.5  # heading + spacing
    elif btype == "heading4":
        return 2.0
    return 1.0


def _total_visual_lines(content):
    """Sum visual lines for a list of content blocks."""
    return sum(_visual_lines(b) for b in content)


# Max visual lines per slide (code-line equivalents)
MAX_SLIDE_LINES = 18


def parse_html_to_slides(html_path):
    """Parse an HTML file and return a list of slide dicts."""
    with open(html_path, "r", encoding="utf-8") as f:
        soup = BeautifulSoup(f, "html.parser")

    title_tag = soup.find("title")
    title = title_tag.get_text() if title_tag else Path(html_path).stem

    header = soup.find("header")
    if header:
        header.decompose()

    slides = []
    current_heading = None
    current_content = []

    def flush_slide():
        nonlocal current_heading, current_content
        if current_heading is not None:
            slides.append({
                "heading": current_heading,
                "content": current_content,
            })
        current_heading = None
        current_content = []

    body = soup.find("body")
    if not body:
        return title, slides

    for element in body.children:
        if isinstance(element, NavigableString):
            continue
        if not isinstance(element, Tag):
            continue

        tag_name = element.name

        if tag_name in ("h1", "h2"):
            flush_slide()
            current_heading = extract_text(element)
            current_content = []

        elif tag_name == "h3":
            if current_heading is None:
                flush_slide()
                current_heading = extract_text(element)
                current_content = []
            else:
                vlines = _total_visual_lines(current_content)
                if vlines > MAX_SLIDE_LINES * 0.6:
                    # h3 is a natural break point — split early
                    flush_slide()
                    current_heading = extract_text(element)
                    current_content = []
                else:
                    current_content.append({
                        "type": "heading3",
                        "text": extract_text(element),
                    })

        elif tag_name == "h4":
            if current_heading is None:
                flush_slide()
                current_heading = extract_text(element)
                current_content = []
            else:
                current_content.append({
                    "type": "heading4",
                    "text": extract_text(element),
                })

        elif tag_name == "p":
            if current_heading is None:
                current_heading = title
            current_content.append({
                "type": "paragraph",
                "text": extract_text(element),
            })

        elif tag_name in ("ul", "ol"):
            if current_heading is None:
                current_heading = title
            for li in element.find_all("li", recursive=False):
                current_content.append({
                    "type": "list_item",
                    "text": extract_text(li),
                    "level": 0,
                })
                for sub_li in li.find_all("li"):
                    if sub_li.parent != li:
                        continue
                    current_content.append({
                        "type": "list_item",
                        "text": extract_text(sub_li),
                        "level": 1,
                    })

        elif tag_name == "pre" or (tag_name == "div" and "sourceCode" in element.get("class", [])):
            if current_heading is None:
                current_heading = title

            # Try to extract syntax-highlighted tokens
            code_lines = extract_code_tokens(element)
            if code_lines:
                # Cap very long code blocks
                if len(code_lines) > 40:
                    code_lines = code_lines[:40]
                    code_lines.append([("... (continued)", None)])
                current_content.append({
                    "type": "code_highlighted",
                    "text": "",  # not used for rendering but needed for length calc
                    "lines": code_lines,
                })
            else:
                code_text = extract_code_text(element)
                if len(code_text) > 2000:
                    code_text = code_text[:2000] + "\n... (continued)"
                current_content.append({
                    "type": "code",
                    "text": code_text,
                })

        elif tag_name == "table":
            if current_heading is None:
                current_heading = title
            current_content.append({
                "type": "table_text",
                "text": table_to_text(element),
            })

        elif tag_name == "blockquote":
            if current_heading is None:
                current_heading = title
            current_content.append({
                "type": "blockquote",
                "text": extract_text(element),
            })

        elif tag_name == "hr":
            if current_content:
                vlines = _total_visual_lines(current_content)
                if vlines > MAX_SLIDE_LINES * 0.4:
                    flush_slide()

    flush_slide()

    # Split slides that exceed visual line budget.
    # First, break apart oversized code blocks so they can be distributed.
    def _split_large_blocks(content):
        """Split code blocks that are taller than a full slide."""
        result = []
        for block in content:
            btype = block["type"]
            if btype == "code_highlighted":
                lines = block["lines"]
                if len(lines) <= MAX_SLIDE_LINES:
                    result.append(block)
                else:
                    for start in range(0, len(lines), MAX_SLIDE_LINES):
                        result.append({
                            "type": "code_highlighted",
                            "text": "",
                            "lines": lines[start:start + MAX_SLIDE_LINES],
                        })
            elif btype == "code":
                code_lines = block["text"].split("\n")
                if len(code_lines) <= MAX_SLIDE_LINES:
                    result.append(block)
                else:
                    for start in range(0, len(code_lines), MAX_SLIDE_LINES):
                        result.append({
                            "type": "code",
                            "text": "\n".join(code_lines[start:start + MAX_SLIDE_LINES]),
                        })
            elif btype == "table_text":
                tbl_lines = block["text"].split("\n")
                if len(tbl_lines) <= MAX_SLIDE_LINES:
                    result.append(block)
                else:
                    for start in range(0, len(tbl_lines), MAX_SLIDE_LINES):
                        result.append({
                            "type": "table_text",
                            "text": "\n".join(tbl_lines[start:start + MAX_SLIDE_LINES]),
                        })
            else:
                result.append(block)
        return result

    final_slides = []
    for slide_data in slides:
        content = _split_large_blocks(slide_data["content"])
        heading = slide_data["heading"]

        total_vl = _total_visual_lines(content)
        if total_vl <= MAX_SLIDE_LINES:
            final_slides.append({"heading": heading, "content": content})
            continue

        # Greedily pack blocks into slides
        chunk = []
        chunk_vl = 0.0
        part = 1
        for block in content:
            bvl = _visual_lines(block)
            if chunk_vl + bvl > MAX_SLIDE_LINES and chunk:
                final_slides.append({
                    "heading": f"{heading} (cont.)" if part > 1 else heading,
                    "content": chunk,
                })
                chunk = []
                chunk_vl = 0.0
                part += 1
            chunk.append(block)
            chunk_vl += bvl
        if chunk:
            final_slides.append({
                "heading": f"{heading} (cont.)" if part > 1 else heading,
                "content": chunk,
            })

    return title, final_slides


def add_section_divider(prs, title):
    """Add a section divider slide between guides."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    left = Inches(1)
    top = Inches(2.8)
    width = Inches(11.333)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = SPRYKER_TEAL
    p.alignment = PP_ALIGN.CENTER

    line_shape = slide.shapes.add_shape(
        1, Inches(4), Inches(2.5), Inches(5.333), Pt(3)
    )
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = SPRYKER_TEAL
    line_shape.line.fill.background()

    return slide


def build_combined_pptx(html_files, output_path, pptx_title, subtitle, objectives):
    """Combine multiple HTML guides into a single PowerPoint file."""
    prs = create_presentation()

    # Title slide
    add_title_slide(prs, pptx_title, subtitle)
    total_slides = 1

    # Objectives slide
    add_objectives_slide(prs, objectives)
    total_slides += 1

    for html_path in html_files:
        title, slides = parse_html_to_slides(str(html_path))

        add_section_divider(prs, title)
        total_slides += 1

        for slide_data in slides:
            add_content_slide(prs, slide_data["heading"], slide_data["content"])
            total_slides += 1

        print(f"  + {html_path.name} ({len(slides)} content slides)")

    prs.save(output_path)
    return total_slides


def main():
    base_dir = Path(__file__).parent / "guides-html"
    output_dir = Path(__file__).parent / "guides-pptx"
    output_dir.mkdir(exist_ok=True)

    configs = [
        {
            "subdir": "basics",
            "title": "Spryker Academy \u2014 Fundamental Exercises",
            "subtitle": "Instructor-Led Training | Backend Development Fundamentals",
            "output": "spryker-academy-fundamental.pptx",
            "objectives": BASICS_OBJECTIVES,
        },
        {
            "subdir": "intermediate",
            "title": "Spryker Academy \u2014 Professional Exercises",
            "subtitle": "Instructor-Led Training | Advanced Spryker Development",
            "output": "spryker-academy-professional.pptx",
            "objectives": INTERMEDIATE_OBJECTIVES,
        },
    ]

    for cfg in configs:
        dir_path = base_dir / cfg["subdir"]
        html_files = sorted(dir_path.glob("*.html"))
        if not html_files:
            print(f"No HTML files in {cfg['subdir']}/")
            continue

        out_file = output_dir / cfg["output"]
        print(f"\n{cfg['title']} ({len(html_files)} guides):")
        total = build_combined_pptx(
            html_files, str(out_file), cfg["title"], cfg["subtitle"], cfg["objectives"]
        )
        print(f"  => {out_file.name} \u2014 {total} slides total")

    print(f"\nDone! Output in: {output_dir}")


if __name__ == "__main__":
    main()
